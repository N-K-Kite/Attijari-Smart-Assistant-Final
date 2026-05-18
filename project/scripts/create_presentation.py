import os
import glob
from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

title_slide_layout = prs.slide_layouts[0]
title_and_content_layout = prs.slide_layouts[1]

def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text
    subtitle.text = subtitle_text

def add_content_slide(prs, title_text, content_text):
    slide = prs.slides.add_slide(title_and_content_layout)
    title = slide.shapes.title
    title.text = title_text
    content = slide.placeholders[1]
    content.text = content_text

def add_image_slide(prs, title_text, image_path):
    slide = prs.slides.add_slide(title_and_content_layout)
    title = slide.shapes.title
    title.text = title_text
    
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, left, top, width=width)

add_title_slide(prs, "Système de Détection IA & RPA", "Projet de Fin d'Études 2026\nAttijari Bank")

add_content_slide(prs, "Contexte & Problématique", 
                  "• Augmentation des volumes de tickets IT.\n"
                  "• Temps de diagnostic manuel chronophage.\n"
                  "• Risque élevé de rater les incidents critiques parmi les requêtes standard.\n"
                  "• Nécessité d'automatiser le tri et les actions curatives simples.")

add_content_slide(prs, "Objectifs du Système",
                  "• DÉTECTER : Modèle LSTM pour repérer les comportements anormaux sur l'infrastructure.\n"
                  "• ANALYSER : Moteur NLP pour extraire le sens des réclamations.\n"
                  "• RECOMMANDER : Algorithme KNN pour suggérer les solutions historiques.\n"
                  "• AGIR : RPA (UiPath) pour automatiser la résolution sans intervention humaine.")

add_content_slide(prs, "Architecture & Stack Technique",
                  "• Interface : Angular 21 (SPA, Dashboard Réactif).\n"
                  "• Cœur Logique : FastAPI (Asynchrone, API REST, OpenAPI).\n"
                  "• Intelligence : TensorFlow (LSTM), Scikit-Learn (KNN, NLP).\n"
                  "• Stockage : PostgreSQL (Données relationnelles et audit).")

screenshots_dir = r"c:\Users\21655\Desktop\attijari-pfe\screenshots"
screenshots = sorted(glob.glob(os.path.join(screenshots_dir, "*.png")))

titles = [
    "Tableau de Bord Principal (Dashboard)",
    "Vue des Recommandations KNN",
    "Moteur d'Analyse IA Temps Réel",
    "Suivi des Prédictions LSTM",
    "Gestion des Alertes & RPA",
    "Boîte de Réception des Réclamations",
    "Détail d'une Réclamation",
    "Interface Chatbot IA",
    "Traçabilité et Audit Trail"
]

for i, img_path in enumerate(screenshots):
    slide_title = titles[i] if i < len(titles) else f"Démonstration Fonctionnelle {i+1}"
    add_image_slide(prs, slide_title, img_path)

add_content_slide(prs, "Conclusion",
                  "• Intégration complète : De la déclaration au RPA.\n"
                  "• Modèles IA opérationnels avec des scores de confiance de 87%.\n"
                  "• Prochaines étapes : Finalisation des workflows UiPath et passage en recette.\n"
                  "• Merci de votre attention.")

output_path = r"c:\Users\21655\Desktop\attijari-pfe\Presentation_PFE_Attijari.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
